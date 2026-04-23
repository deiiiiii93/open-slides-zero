from __future__ import annotations

import json
from io import BytesIO
from pathlib import Path

import pytest
from fastapi import HTTPException, UploadFile
from fastapi.testclient import TestClient

from app.api import decks
from app.artifacts import store
from app.graph import graph as graph_module
from app.graph.nodes import ingest, outline
from app.llm.models import get_model
from app.main import app


def _write_pdf(path: Path, pages: list[str | None]) -> None:
    def escape_pdf_text(text: str) -> str:
        return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")

    objects: list[str] = []
    page_refs = " ".join(f"{idx + 3} 0 R" for idx in range(len(pages)))
    font_obj_num = len(pages) + 3
    content_obj_start = font_obj_num + 1

    objects.append("<< /Type /Catalog /Pages 2 0 R >>")
    objects.append(f"<< /Type /Pages /Kids [{page_refs}] /Count {len(pages)} >>")

    for idx in range(len(pages)):
        objects.append(
            "<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            f"/Resources << /Font << /F1 {font_obj_num} 0 R >> >> "
            f"/Contents {content_obj_start + idx} 0 R >>"
        )
    objects.append("<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    for text in pages:
        stream = ""
        if text is not None:
            stream = f"BT /F1 18 Tf 72 720 Td ({escape_pdf_text(text)}) Tj ET"
        objects.append(
            f"<< /Length {len(stream.encode('latin-1'))} >>\nstream\n{stream}\nendstream"
        )

    parts: list[bytes] = [b"%PDF-1.4\n"]
    offsets = [0]
    for obj_num, obj in enumerate(objects, start=1):
        offsets.append(sum(len(part) for part in parts))
        parts.append(f"{obj_num} 0 obj\n{obj}\nendobj\n".encode("latin-1"))

    xref_start = sum(len(part) for part in parts)
    xref = [b"xref\n", f"0 {len(objects) + 1}\n".encode("latin-1"), b"0000000000 65535 f \n"]
    for offset in offsets[1:]:
        xref.append(f"{offset:010d} 00000 n \n".encode("latin-1"))
    trailer = (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_start}\n%%EOF\n"
    ).encode("latin-1")
    path.write_bytes(b"".join(parts + xref + [trailer]))


def _parse_sse_events(body: str) -> list[dict]:
    events: list[dict] = []
    for line in body.splitlines():
        if not line.startswith("data:"):
            continue
        events.append(json.loads(line[len("data:"):].strip()))
    return events


@pytest.fixture()
def isolated_graph(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    graph_module._compiled = None
    db_path = tmp_path / "threads.sqlite"
    monkeypatch.setattr(graph_module, "DB_PATH", db_path)
    # decks.py imports DB_PATH by name, so its module-local ref needs patching too.
    monkeypatch.setattr(decks, "DB_PATH", db_path)
    monkeypatch.setattr(store, "ROOT", tmp_path / "threads")

    def fake_chat_structured(_model, _messages, schema, **_kwargs):
        if schema is outline._ProposedChoices:
            return schema(
                recommended_scenario_id="sales_pitch",
                candidate_structure_ids=["pyramid", "bluf"],
                rationale="fit",
            )
        raise AssertionError(f"Unexpected schema: {schema}")

    monkeypatch.setattr(outline.zenmux, "chat_structured", fake_chat_structured)
    yield
    graph_module._compiled = None


@pytest.mark.parametrize("suffix", [".txt", ".md", ".markdown"])
def test_parse_text_and_markdown_files_verbatim(tmp_path: Path, suffix: str):
    path = tmp_path / f"notes{suffix}"
    contents = "Revenue up\nMargin steady"
    path.write_text(contents, encoding="utf-8")

    parsed = ingest._parse_file_material(str(path))

    assert parsed.text == contents
    assert parsed.note is None


def test_parse_docx_extracts_paragraphs_and_tables(tmp_path: Path):
    from docx import Document

    path = tmp_path / "brief.docx"
    document = Document()
    document.add_paragraph("Executive summary")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "120"
    document.save(path)

    parsed = ingest._parse_file_material(str(path))

    assert "Executive summary" in parsed.text
    assert "| Metric | Value |" in parsed.text
    assert "| Revenue | 120 |" in parsed.text


def test_parse_pptx_extracts_slide_text_tables_and_notes(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly update"
    slide.placeholders[1].text = "Growth remains above plan"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1.5)).table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Growth"
    table.cell(1, 0).text = "APAC"
    table.cell(1, 1).text = "18%"
    slide.notes_slide.notes_text_frame.text = "Speak to the margin bridge."
    presentation.save(path)

    parsed = ingest._parse_file_material(str(path))

    assert "Quarterly update" in parsed.text
    assert "Growth remains above plan" in parsed.text
    assert "| Region | Growth |" in parsed.text
    assert "Speaker notes" in parsed.text
    assert "Speak to the margin bridge." in parsed.text


def test_parse_xlsx_extracts_non_empty_range(tmp_path: Path):
    from openpyxl import Workbook

    path = tmp_path / "model.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = "Metric"
    sheet["B1"] = "Value"
    sheet["A2"] = "Revenue"
    sheet["B2"] = 120
    workbook.save(path)

    parsed = ingest._parse_file_material(str(path))

    assert "Sheet: Summary" in parsed.text
    assert "Range: A1:B2" in parsed.text
    assert "| Revenue | 120 |" in parsed.text


def test_parse_pdf_extracts_text_and_marks_single_image_only_page(tmp_path: Path):
    path = tmp_path / "report.pdf"
    _write_pdf(path, ["Revenue increased 20%", None])

    parsed = ingest._parse_file_material(str(path))

    assert "Revenue increased 20%" in parsed.text
    assert parsed.note == "Page 2 had no extractable PDF text and was treated as image-only."


def test_parse_pdf_notes_multiple_image_only_pages(tmp_path: Path):
    path = tmp_path / "report.pdf"
    _write_pdf(path, ["Intro", None, "Conclusion", None])

    parsed = ingest._parse_file_material(str(path))

    assert "Intro" in parsed.text
    assert "Conclusion" in parsed.text
    assert parsed.note == "Pages 2, 4 had no extractable PDF text and were treated as image-only."


def test_extract_notes_text_filters_housekeeping_placeholders():
    from types import SimpleNamespace

    def shape(text: str, *, has_text_frame: bool = True, placeholder_name: str | None = None):
        placeholder_format = (
            SimpleNamespace(type=SimpleNamespace(name=placeholder_name))
            if placeholder_name is not None
            else None
        )
        return SimpleNamespace(
            has_text_frame=has_text_frame,
            text=text,
            placeholder_format=placeholder_format,
        )

    slide = SimpleNamespace(
        notes_slide=SimpleNamespace(
            shapes=[
                shape("Real notes", placeholder_name="BODY"),
                shape("7", placeholder_name="SLIDE_NUMBER"),
                shape("2026-04-23", placeholder_name="DATE"),
                shape("Confidential", placeholder_name="FOOTER"),
                shape("Header line", placeholder_name="HEADER"),
                shape("", has_text_frame=False),  # slide image (no text frame)
                shape("Plain shape, no placeholder", placeholder_name=None),
            ]
        )
    )

    assert ingest._extract_notes_text(slide) == "Real notes\nPlain shape, no placeholder"


def test_image_material_uses_zenmux_ocr_model(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    path = tmp_path / "chart.png"
    path.write_bytes(b"fake-image")
    seen: dict[str, object] = {}

    def fake_chat(model, _messages, **kwargs):
        seen["model"] = model
        seen["images"] = kwargs.get("images")
        return "Revenue chart\nVisual summary\n- Blue bars by region"

    monkeypatch.setattr(ingest.zenmux, "chat", fake_chat)

    parsed = ingest._parse_file_material(str(path))

    assert seen["model"] == get_model("ingest.ocr")
    assert seen["images"] == [str(path)]
    assert "Revenue chart" in parsed.text


def test_ingest_node_raises_when_all_materials_are_empty(tmp_path: Path):
    path = tmp_path / "empty.txt"
    path.write_text("", encoding="utf-8")

    with pytest.raises(ValueError, match="No usable content could be extracted"):
        ingest.ingest_node({"materials": [{"kind": "file", "uri": str(path), "name": "empty.txt"}]})


@pytest.mark.asyncio
async def test_normalize_uploaded_materials_uniquifies_duplicate_filenames(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    monkeypatch.setattr(store, "ROOT", tmp_path / "threads")

    materials = await decks.normalize_uploaded_materials(
        "thread-1",
        [
            UploadFile(filename="dup.txt", file=BytesIO(b"one")),
            UploadFile(filename="dup.txt", file=BytesIO(b"two")),
        ],
    )

    assert [material.name for material in materials] == ["dup.txt", "dup.txt"]
    assert materials[0].uri != materials[1].uri
    assert Path(materials[0].uri).name == "dup.txt"
    assert Path(materials[1].uri).name == "dup-2.txt"


@pytest.mark.asyncio
async def test_normalize_uploaded_materials_rejects_unsupported_extensions(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    monkeypatch.setattr(store, "ROOT", tmp_path / "threads")

    with pytest.raises(HTTPException) as exc_info:
        await decks.normalize_uploaded_materials(
            "thread-1",
            [UploadFile(filename="payload.exe", file=BytesIO(b"bad"))],
        )

    assert exc_info.value.status_code == 400
    assert "Unsupported file type" in exc_info.value.detail


def test_upload_stream_preserves_text_then_file_order(isolated_graph):
    client = TestClient(app)

    response = client.post(
        "/decks/upload/stream",
        data={
            "deck_name": "",
            "text": "Lead with the headline",
            "expected_pages": "6",
            "aspect_ratio": "16:9",
            "density_preference": "balanced",
            "language": "en",
            "visual_style_preference": "Editorial",
        },
        files=[
            ("files", ("alpha.txt", b"Revenue up", "text/plain")),
            ("files", ("beta.markdown", b"# Margin\nStable", "text/markdown")),
        ],
    )

    assert response.status_code == 200
    events = _parse_sse_events(response.text)
    done = next(event for event in events if event["type"] == "done")
    materials = done["state"]["values"]["materials"]

    assert [material["name"] for material in materials if material.get("name")] == ["alpha.txt", "beta.markdown"]
    assert materials[0]["kind"] == "text"
    assert materials[0]["parsed"] == "Lead with the headline"
    assert materials[1]["parsed"] == "Revenue up"
    assert materials[2]["parsed"] == "# Margin\nStable"


def test_upload_stream_returns_error_when_all_sources_are_empty(isolated_graph):
    client = TestClient(app)

    response = client.post(
        "/decks/upload/stream",
        data={
            "expected_pages": "4",
            "aspect_ratio": "16:9",
            "density_preference": "balanced",
            "language": "en",
        },
        files=[("files", ("empty.txt", b"", "text/plain"))],
    )

    assert response.status_code == 200
    events = _parse_sse_events(response.text)
    error = next(event for event in events if event["type"] == "error")
    assert "No usable content could be extracted" in error["message"]

    # Cleanup invariant: the failed thread must not linger in /decks or on disk.
    thread_event = next(event for event in events if event["type"] == "thread")
    thread_id = thread_event["thread_id"]
    decks_response = client.get("/decks")
    thread_ids = [deck["thread_id"] for deck in decks_response.json()["decks"]]
    assert thread_id not in thread_ids
    assert not (store.ROOT / thread_id).exists()


def test_create_deck_rejects_empty_materials(isolated_graph):
    client = TestClient(app)

    response = client.post("/decks", json={"materials": []})

    assert response.status_code == 400
    assert "at least one file or paste text" in response.json()["detail"]
    assert client.get("/decks").json()["decks"] == []


def test_stream_create_deck_rejects_empty_materials(isolated_graph):
    client = TestClient(app)

    response = client.post("/decks/stream", json={"materials": []})

    assert response.status_code == 400
    assert "at least one file or paste text" in response.json()["detail"]
    assert client.get("/decks").json()["decks"] == []


def test_upload_stream_unsupported_extension_leaves_no_zombie(isolated_graph):
    client = TestClient(app)

    response = client.post(
        "/decks/upload/stream",
        data={"expected_pages": "4", "aspect_ratio": "16:9", "density_preference": "balanced", "language": "en"},
        files=[("files", ("payload.exe", b"bad", "application/octet-stream"))],
    )

    assert response.status_code == 400
    assert "Unsupported file type" in response.json()["detail"]
    assert client.get("/decks").json()["decks"] == []


def test_upload_stream_rejects_oversize_single_file(isolated_graph, monkeypatch: pytest.MonkeyPatch):
    monkeypatch.setattr(decks, "MAX_UPLOAD_BYTES", 1024)
    client = TestClient(app)

    oversize_payload = b"A" * (1024 + 1)
    response = client.post(
        "/decks/upload/stream",
        data={"expected_pages": "4", "aspect_ratio": "16:9", "density_preference": "balanced", "language": "en"},
        files=[("files", ("big.txt", oversize_payload, "text/plain"))],
    )

    assert response.status_code == 413
    assert "per-file upload limit" in response.json()["detail"]
    assert client.get("/decks").json()["decks"] == []
    # Cleanup invariant: no orphaned bytes from the partially-streamed file.
    materials_dirs = list((store.ROOT).glob("*/materials")) if store.ROOT.exists() else []
    assert all(not any(p.is_file() for p in d.iterdir()) for d in materials_dirs)


def test_upload_stream_rejects_oversize_request_total(isolated_graph, monkeypatch: pytest.MonkeyPatch):
    monkeypatch.setattr(decks, "MAX_UPLOAD_BYTES", 10_000)
    monkeypatch.setattr(decks, "MAX_REQUEST_BYTES", 1500)
    client = TestClient(app)

    response = client.post(
        "/decks/upload/stream",
        data={"expected_pages": "4", "aspect_ratio": "16:9", "density_preference": "balanced", "language": "en"},
        files=[
            ("files", ("one.txt", b"A" * 800, "text/plain")),
            ("files", ("two.txt", b"B" * 800, "text/plain")),
        ],
    )

    assert response.status_code == 413
    assert "per-request limit" in response.json()["detail"]
    assert client.get("/decks").json()["decks"] == []
    # Any already-stored file from the first upload should have been swept.
    assert not store.ROOT.exists() or not any(store.ROOT.iterdir())
