"""Stage B — parse uploaded materials into normalized text + metadata.

Supported inputs:
  - inline / local text: txt, md, markdown
  - office docs: docx, pptx, xlsx
  - pdf: selectable text when present; image-only pages fall back to OCR
  - raster images: OCR (glm-ocr when OSZ_MODEL_OCR_* is set, else ZenMux)
"""

from __future__ import annotations

import base64
import io
import os
from dataclasses import dataclass
import logging
from pathlib import Path
from typing import Any

from ...llm import glm_ocr, zenmux
from ...llm.models import get_model

log = logging.getLogger(__name__)

TEXT_SUFFIXES = {".txt", ".md", ".markdown"}
IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png"}
_NOTES_SKIP_PLACEHOLDERS = frozenset({"SLIDE_NUMBER", "DATE", "FOOTER", "HEADER"})

# Defensive cap on OCR calls for a single PDF. A 25MB scan can hold hundreds of
# pages; without this cap a single upload could spray the OCR model. Pages past
# the cap are listed in the note so the user knows what we skipped.
_PDF_OCR_MAX_PAGES = int(os.getenv("OSZ_PDF_OCR_MAX_PAGES", "40"))
# 2x scale ≈ 144 DPI — legible OCR input without bloating the base64 payload.
_PDF_OCR_RENDER_SCALE = float(os.getenv("OSZ_PDF_OCR_RENDER_SCALE", "2.0"))


@dataclass
class ParsedMaterial:
    text: str
    note: str | None = None


def _append_note(existing: str | None, extra: str | None) -> str | None:
    if existing and extra:
        return f"{existing}\n{extra}"
    return extra or existing


def _escape_table_cell(value: str) -> str:
    return value.replace("|", r"\|").replace("\n", "<br>")


def _markdown_table(rows: list[list[str]]) -> str:
    normalized = [
        [cell.strip() for cell in row]
        for row in rows
        if any(cell.strip() for cell in row)
    ]
    if not normalized:
        return ""
    width = max(len(row) for row in normalized)
    padded = [row + [""] * (width - len(row)) for row in normalized]
    header = padded[0]
    separator = ["---"] * width
    lines = [
        "| " + " | ".join(_escape_table_cell(cell) for cell in header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in padded[1:]:
        lines.append("| " + " | ".join(_escape_table_cell(cell) for cell in row) + " |")
    return "\n".join(lines)


def _finalize_text(text: str, note: str | None = None) -> ParsedMaterial:
    normalized = text.strip()
    if normalized:
        return ParsedMaterial(text=normalized, note=note)
    return ParsedMaterial(text="", note=note or "No extractable content found.")


_OCR_PROMPT = (
    "Extract slide-source material from this image using OCR.\n"
    "1. Transcribe visible text faithfully, preserving the original language.\n"
    "2. Reconstruct tables as Markdown when possible.\n"
    "3. Include chart titles, axes labels, legends, units, and notable values.\n"
    "4. If the image is mostly non-text, append a short 'Visual summary' section "
    "covering the key visual evidence.\n"
    "Return plain text only."
)

_IMAGE_MIME_BY_SUFFIX = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg"}


def _to_image_data_uri(uri: str) -> str:
    """Normalize an image reference into something glm-ocr accepts: an
    http(s) URL, or an image data URI built from a local path.
    Pass-through if already a URL or data URI."""
    if uri.startswith(("http://", "https://", "data:")):
        return uri
    path = Path(uri)
    mime = _IMAGE_MIME_BY_SUFFIX.get(path.suffix.lower(), "image/png")
    encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def _ocr_via_zenmux(image_uri: str) -> str:
    raw = zenmux.chat(
        get_model("ingest.ocr"),
        [{"role": "user", "content": _OCR_PROMPT}],
        images=[image_uri],
        temperature=0.1,
    )
    return (raw or "").strip()


def _ocr_image_source(image_uri: str) -> str:
    """Run OCR on a local path, data URI, or http(s) URL. Returns stripped text.

    Routes to the standalone GLM-OCR endpoint when the OSZ_MODEL_OCR_* env
    vars are set; falls back to the ZenMux vision path otherwise.
    """
    if glm_ocr.is_configured():
        return glm_ocr.extract_markdown(_to_image_data_uri(image_uri))
    return _ocr_via_zenmux(image_uri)


def _parse_image_ocr(uri: str) -> ParsedMaterial:
    return _finalize_text(_ocr_image_source(uri))


def _read_text_file(uri: str) -> ParsedMaterial:
    if uri.startswith("text:"):
        return _finalize_text(uri[len("text:"):], note=None)
    path = Path(uri)
    if not path.exists():
        raise FileNotFoundError(f"Material file not found: {uri}")
    return _finalize_text(path.read_text(encoding="utf-8", errors="replace"))


def _pages_phrase(pages: list[int]) -> str:
    if len(pages) == 1:
        return f"Page {pages[0]}"
    return "Pages " + ", ".join(str(p) for p in pages)


def _open_pdf(uri: str):
    """Open a PDF via pypdfium2, retrying with an empty password on failure.

    pypdfium2 wraps PDFium, which handles AES/RC4 encryption natively — unlike
    pypdf, which needs the `cryptography` dependency for AES. Many "protected"
    PDFs from firms (KPMG, E&Y, government portals) use owner-password
    encryption with an empty user password; retrying with password="" unlocks
    those without prompting.
    """
    import pypdfium2 as pdfium

    try:
        return pdfium.PdfDocument(uri)
    except pdfium.PdfiumError:
        return pdfium.PdfDocument(uri, password="")


def _extract_page_text(page) -> str:
    try:
        return (page.get_textpage().get_text_range() or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def _render_page_to_data_uri(page) -> str:
    """Rasterize a pypdfium2 page to a PNG data URI."""
    bitmap = page.render(scale=_PDF_OCR_RENDER_SCALE)
    image = bitmap.to_pil()
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/png;base64,{encoded}"


def _parse_pdf(uri: str) -> ParsedMaterial:
    document = _open_pdf(uri)
    page_blocks: dict[int, str] = {}
    image_only_pages: list[int] = []
    pages_by_number: dict[int, Any] = {}
    try:
        for idx, page in enumerate(document, start=1):
            pages_by_number[idx] = page
            text = _extract_page_text(page)
            if text:
                page_blocks[idx] = f"## Page {idx}\n{text}"
            else:
                image_only_pages.append(idx)

        note_fragments: list[str] = []
        if image_only_pages:
            ocr_pages = image_only_pages[:_PDF_OCR_MAX_PAGES]
            skipped_pages = image_only_pages[_PDF_OCR_MAX_PAGES:]
            ocr_recovered: list[int] = []
            ocr_failed: list[int] = []
            for page_number in ocr_pages:
                try:
                    data_uri = _render_page_to_data_uri(pages_by_number[page_number])
                    ocr_text = _ocr_image_source(data_uri)
                except Exception as exc:  # noqa: BLE001
                    log.warning("OCR failed for PDF page %d: %s", page_number, exc)
                    ocr_failed.append(page_number)
                    continue
                if ocr_text:
                    page_blocks[page_number] = f"## Page {page_number} (OCR)\n{ocr_text}"
                    ocr_recovered.append(page_number)
                else:
                    ocr_failed.append(page_number)

            if ocr_recovered:
                note_fragments.append(
                    f"{_pages_phrase(ocr_recovered)} had no embedded PDF text; "
                    "recovered via OCR."
                )
            if ocr_failed:
                note_fragments.append(
                    f"{_pages_phrase(ocr_failed)} had no extractable text; "
                    "OCR could not recover content."
                )
            if skipped_pages:
                note_fragments.append(
                    f"{_pages_phrase(skipped_pages)} exceeded the OCR page cap "
                    f"({_PDF_OCR_MAX_PAGES}) and were skipped."
                )
    finally:
        document.close()

    ordered = [page_blocks[i] for i in sorted(page_blocks)]
    note = "\n".join(note_fragments) if note_fragments else None
    return _finalize_text("\n\n".join(ordered), note=note)


def _parse_docx(uri: str) -> ParsedMaterial:
    from docx import Document

    document = Document(uri)
    blocks: list[str] = []

    paragraphs = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]
    if paragraphs:
        blocks.append("\n".join(paragraphs))

    for idx, table in enumerate(document.tables, start=1):
        rows = [
            [cell.text.strip() for cell in row.cells]
            for row in table.rows
        ]
        table_md = _markdown_table(rows)
        if table_md:
            blocks.append(f"## Table {idx}\n{table_md}")

    return _finalize_text("\n\n".join(blocks))


def _extract_notes_text(slide: Any) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    texts: list[str] = []
    for shape in getattr(notes_slide, "shapes", []):
        if not getattr(shape, "has_text_frame", False):
            continue
        ph = getattr(shape, "placeholder_format", None)
        ph_type = getattr(ph, "type", None) if ph is not None else None
        ph_name = getattr(ph_type, "name", None) if ph_type is not None else None
        if ph_name in _NOTES_SKIP_PLACEHOLDERS:
            continue
        text = (getattr(shape, "text", "") or "").strip()
        if text:
            texts.append(text)
    return "\n".join(texts)


def _parse_pptx(uri: str) -> ParsedMaterial:
    from pptx import Presentation

    presentation = Presentation(uri)
    slides: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = getattr(shape, "text", "").strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                table_md = _markdown_table(rows)
                if table_md:
                    parts.append(f"Table\n{table_md}")
        notes_text = _extract_notes_text(slide)
        if notes_text:
            parts.append(f"Speaker notes\n{notes_text}")
        if parts:
            slides.append(f"## Slide {idx}\n" + "\n\n".join(parts))

    return _finalize_text("\n\n".join(slides))


def _cell_to_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _parse_xlsx(uri: str) -> ParsedMaterial:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    workbook = load_workbook(uri, data_only=True, read_only=True)
    blocks: list[str] = []

    for sheet in workbook.worksheets:
        all_rows = list(sheet.iter_rows(values_only=True))
        min_row: int | None = None
        max_row = 0
        min_col: int | None = None
        max_col = 0
        for row_idx, row in enumerate(all_rows, start=1):
            for col_idx, value in enumerate(row, start=1):
                if _cell_to_text(value):
                    min_row = row_idx if min_row is None else min(min_row, row_idx)
                    max_row = max(max_row, row_idx)
                    min_col = col_idx if min_col is None else min(min_col, col_idx)
                    max_col = max(max_col, col_idx)
        if min_row is None or min_col is None:
            continue

        extracted_rows: list[list[str]] = []
        for row_idx in range(min_row, max_row + 1):
            source_row = all_rows[row_idx - 1] if row_idx - 1 < len(all_rows) else ()
            extracted_rows.append(
                [
                    _cell_to_text(source_row[col_idx - 1] if col_idx - 1 < len(source_row) else None)
                    for col_idx in range(min_col, max_col + 1)
                ]
            )
        table_md = _markdown_table(extracted_rows)
        if table_md:
            coord = f"{get_column_letter(min_col)}{min_row}:{get_column_letter(max_col)}{max_row}"
            blocks.append(f"## Sheet: {sheet.title}\nRange: {coord}\n{table_md}")

    return _finalize_text("\n\n".join(blocks))


def _parse_file_material(uri: str) -> ParsedMaterial:
    if uri.startswith("text:"):
        return _read_text_file(uri)
    path = Path(uri)
    if not path.exists():
        raise FileNotFoundError(f"Material file not found: {uri}")
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return _read_text_file(uri)
    if suffix in IMAGE_SUFFIXES:
        return _parse_image_ocr(uri)
    if suffix == ".pdf":
        return _parse_pdf(uri)
    if suffix == ".docx":
        return _parse_docx(uri)
    if suffix == ".pptx":
        return _parse_pptx(uri)
    if suffix == ".xlsx":
        return _parse_xlsx(uri)
    raise ValueError(f"Unsupported material type: {suffix or 'unknown'}")


def ingest_node(state: dict[str, Any]) -> dict[str, Any]:
    materials = state.get("materials", [])
    out: list[dict[str, Any]] = []
    for material in materials:
        if material.get("parsed"):
            out.append(material)
            continue
        try:
            parsed = _parse_file_material(material["uri"])
            out.append(
                {
                    **material,
                    "parsed": parsed.text,
                    "note": _append_note(material.get("note"), parsed.note),
                }
            )
        except Exception as exc:  # noqa: BLE001
            log.exception("Failed to parse material %s", material.get("uri"))
            out.append(
                {
                    **material,
                    "parsed": "",
                    "note": _append_note(material.get("note"), f"Failed to parse: {exc}"),
                }
            )

    if not any((material.get("parsed") or "").strip() for material in out):
        raise ValueError("No usable content could be extracted from the provided materials.")

    return {
        "materials": out,
        "current_stage": "outline",
    }
